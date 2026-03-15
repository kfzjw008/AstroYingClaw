#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
AstroGuide Hardware Skill v2.1 - OpenClaw 硬件技能执行脚本
用于控制智慧天文指星仪设备（连续多步讲解 + 课件自动生成版）

被 OpenClaw Agent 通过以下命令调用：
python run_hardware.py '{lesson_script}' '{lecture_md}' '{ppt_slides}'

功能：
1. 接收讲课脚本数组（包含多个步骤的解说词和坐标）
2. 接收长篇讲稿（Markdown 格式，约 2000 字）
3. 接收 PPT 大纲（4-5 页幻灯片）
4. 自动生成本地 Markdown 讲稿和 PowerPoint 课件文件
5. 按顺序遍历每个步骤：发送坐标到串口 -> 播放 TTS 语音 -> 等待完成 -> 下一步
6. 返回 JSON 格式的执行结果给 OpenClaw 框架

作者：AstroClaw Team
许可：MIT License
版本：2.1.1 (UTF-8 编码修复版)
"""

import sys
import os
import io

# ============================================
# UTF-8 编码修复（必须在所有其他导入之前执行）
# ============================================
if sys.platform == 'win32':
    # Windows 下强制设置控制台为 UTF-8 编码
    # 方法 1: 设置环境变量（影响新创建的流）
    os.environ['PYTHONIOENCODING'] = 'utf-8'

    # 方法 2: 重新配置标准流的编码
    # 注意：必须在 colorama 导入之前完成
    try:
        if hasattr(sys.stdout, 'buffer'):
            sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace', newline='\n', write_through=True)
        if hasattr(sys.stderr, 'buffer'):
            sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8', errors='replace', newline='\n', write_through=True)
    except Exception:
        # 如果重新配置失败，继续执行（可能终端本身支持 UTF-8）
        pass

import json
import argparse
import logging
import threading
import queue
import time
from typing import Optional, Dict, Any, List
from datetime import datetime
from pathlib import Path

# ============================================
# 依赖导入检查
# ============================================

# pyserial - 串口通信
try:
    import serial
    import serial.tools.list_ports
    SERIAL_AVAILABLE = True
except ImportError:
    SERIAL_AVAILABLE = False
    print("⚠️ 警告: pyserial 未安装，串口功能将不可用", file=sys.stderr)

# edge-tts - 在线 TTS（推荐）
try:
    import edge_tts
    EDGE_TTS_AVAILABLE = True
except ImportError:
    EDGE_TTS_AVAILABLE = False

# pyttsx3 - 离线 TTS（备选）
try:
    import pyttsx3
    PYTTSX3_AVAILABLE = True
except ImportError:
    PYTTSX3_AVAILABLE = False

# python-pptx - PowerPoint 生成库
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("⚠️ 警告: python-pptx 未安装，PPT 生成功能将不可用", file=sys.stderr)

# ============================================
# 日志配置
# ============================================

class ColoredFormatter(logging.Formatter):
    """带颜色的日志格式化器"""
    COLORS = {
        'DEBUG': '\033[36m',    # 青色
        'INFO': '\033[32m',     # 绿色
        'WARNING': '\033[33m',  # 黄色
        'ERROR': '\033[31m',    # 红色
        'CRITICAL': '\033[35m', # 紫色
    }
    RESET = '\033[0m'

    def format(self, record):
        color = self.COLORS.get(record.levelname, '')
        record.levelname = f"{color}{record.levelname}{self.RESET}"
        return super().format(record)

# 配置日志
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s | %(levelname)-8s | %(message)s',
    datefmt='%H:%M:%S'
)
logger = logging.getLogger(__name__)

# Windows 终端颜色支持
if sys.platform == 'win32':
    try:
        import colorama
        # 使用 autoreset=True 但禁用 convert 和 wrap 来避免 UTF-8 编码问题
        # 只在新版 Windows 10+ 终端支持原生 ANSI 颜色时才启用
        colorama.init(autoreset=True, strip=False, convert=False, wrap=False)
    except (ImportError, Exception):
        # colorama 初始化失败时静默处理，不影响主功能
        pass

# ============================================
# 配置常量
# ============================================

# 串口配置
DEFAULT_BAUDRATE = 115200
DEFAULT_TIMEOUT = 1
SERIAL_PORT_ENV = "ASTRO_GUIDE_SERIAL_PORT"  # 从环境变量读取串口

# TTS 配置
TTS_ENGINE_ENV = "ASTRO_GUIDE_TTS_ENGINE"  # edge 或 pyttsx3
DEFAULT_TTS_ENGINE = "edge"
DEFAULT_VOICE = "zh-CN-XiaoxiaoNeural"  # 默认中文语音

# 步骤间延迟（秒）
STEP_DELAY_ENV = "ASTRO_GUIDE_STEP_DELAY"
DEFAULT_STEP_DELAY = 0.5

# 输出目录配置
OUTPUT_DIR_ENV = "ASTRO_GUIDE_OUTPUT_DIR"
DEFAULT_OUTPUT_DIR = "output"

# Mock 模式标记
FORCE_MOCK_MODE_ENV = "ASTRO_GUIDE_FORCE_MOCK"  # 设为 "true" 强制使用 Mock 模式


# ============================================
# 课件生成器类
# ============================================

class CoursewareGenerator:
    """
    课件自动生成器

    功能：
    - 生成本地 Markdown 讲稿文件
    - 生成 PowerPoint 课件文件
    """

    def __init__(self, output_dir: str = None):
        """
        初始化课件生成器

        Args:
            output_dir (str, optional): 输出目录路径
        """
        self.output_dir = Path(output_dir or os.getenv(OUTPUT_DIR_ENV, DEFAULT_OUTPUT_DIR))

        # 确保输出目录存在
        self.output_dir.mkdir(parents=True, exist_ok=True)
        logger.info(f"📁 输出目录已准备: {self.output_dir}")

    def generate_markdown(self, lecture_md: str, topic: str = "天文讲稿") -> str:
        """
        生成 Markdown 讲稿文件

        Args:
            lecture_md (str): 讲稿内容（Markdown 格式）
            topic (str, optional): 讲稿主题，用于文件名

        Returns:
            str: 生成的文件路径
        """
        try:
            # 生成文件名（带时间戳）
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"{topic}_{timestamp}.md"
            file_path = self.output_dir / filename

            # 写入文件
            with open(file_path, 'w', encoding='utf-8') as f:
                f.write(lecture_md)

            logger.info(f"📝 Markdown 讲稿已生成: {file_path}")
            logger.info(f"   字数: {len(lecture_md)} 字")
            return str(file_path)

        except Exception as e:
            logger.error(f"❌ 生成 Markdown 讲稿失败: {type(e).__name__}: {str(e)}")
            return ""

    def generate_powerpoint(self, ppt_slides: List[Dict[str, Any]], topic: str = "天文课件") -> str:
        """
        生成 PowerPoint 课件文件

        Args:
            ppt_slides (List[Dict]): PPT 页面数组，每页包含 title 和 bullets
            topic (str, optional): 课件主题，用于文件名

        Returns:
            str: 生成的文件路径
        """
        if not PPTX_AVAILABLE:
            logger.error("❌ python-pptx 未安装，无法生成 PPT 文件")
            return ""

        try:
            # 创建 Presentation 对象
            prs = Presentation()

            # 设置幻灯片尺寸（16:9）
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(5.625)

            # 遍历每个幻灯片数据
            for slide_data in ppt_slides:
                title = slide_data.get("title", "")
                bullets = slide_data.get("bullets", [])

                # 添加新幻灯片（使用标题和内容布局）
                slide_layout = prs.slide_layouts[1]  # 1 = Title and Content
                slide = prs.slides.add_slide(slide_layout)

                # 设置标题
                title_shape = slide.shapes.title
                title_shape.text = title

                # 设置内容（项目符号列表）
                content_shape = slide.placeholders[1]  # 1 = Content placeholder
                text_frame = content_shape.text_frame
                text_frame.clear()  # 清除默认内容

                # 添加每个要点
                for idx, bullet in enumerate(bullets):
                    if idx == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()

                    p.text = bullet
                    p.level = 0
                    p.font.size = Pt(18)
                    p.space_after = Pt(12)

            # 生成文件名（带时间戳）
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"{topic}_{timestamp}.pptx"
            file_path = self.output_dir / filename

            # 保存文件
            prs.save(file_path)

            logger.info(f"📊 PowerPoint 课件已生成: {file_path}")
            logger.info(f"   页数: {len(ppt_slides)} 页")
            return str(file_path)

        except Exception as e:
            logger.error(f"❌ 生成 PowerPoint 课件失败: {type(e).__name__}: {str(e)}")
            return ""


# ============================================
# 硬件控制器类
# ============================================

class AstroHardwareController:
    """
    天文指星仪硬件控制器（连续讲解版）

    功能：
    - 管理串口连接（自动降级到 Mock 模式）
    - 发送坐标数据到硬件设备
    - 同步播放 TTS 语音（阻塞等待完成）
    - 支持连续多步讲解
    """

    def __init__(self, serial_port: Optional[str] = None):
        """
        初始化硬件控制器

        Args:
            serial_port (str, optional): 串口设备名称。如未提供，尝试从环境变量读取
        """
        # 确定串口配置
        self.serial_port = serial_port or os.getenv(SERIAL_PORT_ENV, "")
        self.baudrate = DEFAULT_BAUDRATE

        # 串口连接对象
        self.serial_conn: Optional[serial.Serial] = None

        # 工作模式：real 或 mock
        self.mode = "real"  # 默认尝试真实连接

        # TTS 引擎
        self.tts_engine = os.getenv(TTS_ENGINE_ENV, DEFAULT_TTS_ENGINE).lower()

        # 步骤间延迟
        try:
            self.step_delay = float(os.getenv(STEP_DELAY_ENV, DEFAULT_STEP_DELAY))
        except ValueError:
            self.step_delay = DEFAULT_STEP_DELAY

        # 检查是否强制 Mock 模式
        if os.getenv(FORCE_MOCK_MODE_ENV, "").lower() == "true":
            self.mode = "mock"
            logger.info("🎭 强制 Mock 模式（环境变量设置）")
            return

        # 尝试连接串口
        if self.serial_port:
            self._connect_serial()
        else:
            self._try_auto_detect_serial()

    def _try_auto_detect_serial(self):
        """尝试自动检测并连接串口"""
        if not SERIAL_AVAILABLE:
            logger.info("🎭 pyserial 未安装，切换到 Mock 模式")
            self.mode = "mock"
            return

        ports = serial.tools.list_ports.comports()
        if not ports:
            logger.info("🎭 未检测到可用串口，切换到 Mock 模式")
            self.mode = "mock"
            return

        # 尝试连接第一个可用串口
        for port_info in ports:
            logger.info(f"🔍 检测到串口: {port_info.device} - {port_info.description}")
            if self._connect_serial(port_info.device):
                return

        # 所有串口连接失败，使用 Mock 模式
        logger.info("🎭 所有串口连接失败，切换到 Mock 模式")
        self.mode = "mock"

    def _connect_serial(self, port: str = None) -> bool:
        """
        连接串口设备

        Args:
            port (str, optional): 串口名称

        Returns:
            bool: 连接成功返回 True
        """
        target_port = port or self.serial_port

        if not target_port:
            return False

        try:
            if self.serial_conn and self.serial_conn.is_open:
                self.serial_conn.close()

            logger.info(f"🔌 正在连接串口 {target_port} (波特率: {self.baudrate})...")
            self.serial_conn = serial.Serial(
                port=target_port,
                baudrate=self.baudrate,
                timeout=DEFAULT_TIMEOUT,
                write_timeout=DEFAULT_TIMEOUT
            )

            if self.serial_conn.is_open:
                logger.info(f"✅ 串口连接成功 | {target_port}")
                self.mode = "real"
                self.serial_port = target_port
                return True

        except serial.SerialException as e:
            logger.warning(f"⚠️ 串口连接失败: {str(e)}")
        except Exception as e:
            logger.warning(f"⚠️ 连接异常: {type(e).__name__}: {str(e)}")

        return False

    def send_coordinates(self, ra: float, dec: float) -> bool:
        """
        发送天体坐标到硬件设备

        Args:
            ra (float): 赤经（小时），范围 0-24
            dec (float): 赤纬（度），范围 -90~90

        Returns:
            bool: 发送成功返回 True
        """
        # 构造 JSON 数据包
        data_packet = {
            "cmd": "point_to",
            "ra": float(ra),
            "dec": float(dec),
            "timestamp": datetime.now().isoformat()
        }

        json_str = json.dumps(data_packet, ensure_ascii=False)

        if self.mode == "real":
            return self._send_serial(json_str, data_packet)
        else:
            return self._send_mock(json_str, data_packet)

    def _send_serial(self, json_str: str, data_packet: Dict[str, Any]) -> bool:
        """通过真实串口发送数据"""
        if not self.serial_conn or not self.serial_conn.is_open:
            logger.warning("⚠️ 串口未连接，尝试重新连接...")
            if not self._connect_serial():
                logger.info("🔄 重新连接失败，切换到 Mock 模式")
                self.mode = "mock"
                return self._send_mock(json_str, data_packet)

        try:
            message = (json_str + "\n").encode('utf-8')
            bytes_written = self.serial_conn.write(message)
            self.serial_conn.flush()

            logger.info(
                f"📤 [串口] 已发送坐标 | "
                f"RA={data_packet['ra']:.4f}, DEC={data_packet['dec']:.4f} | "
                f"大小={bytes_written} 字节"
            )
            return True

        except Exception as e:
            logger.error(f"❌ 串口发送失败: {type(e).__name__}: {str(e)}")
            self.mode = "mock"
            return self._send_mock(json_str, data_packet)

    def _send_mock(self, json_str: str, data_packet: Dict[str, Any]) -> bool:
        """Mock 模式：打印高亮虚拟日志"""
        # 使用 ANSI 颜色代码高亮显示
        cyan = '\033[36m'
        yellow = '\033[33m'
        green = '\033[32m'
        bold = '\033[1m'
        reset = '\033[0m'

        print()
        print(f"{bold}{cyan}{'='*60}{reset}")
        print(f"{bold}{yellow}🎭 [Mock 模式] 虚拟硬件控制{reset}")
        print(f"{bold}{cyan}{'='*60}{reset}")
        print(f"{green}📡 模拟发送数据到串口:{reset}")
        print(f"   {cyan}JSON:{reset} {json_str}")
        print(f"   {cyan}RA  :{reset} {yellow}{data_packet['ra']:.4f}{reset} 小时")
        print(f"   {cyan}DEC :{reset} {yellow}{data_packet['dec']:.4f}{reset} 度")
        print(f"{bold}{cyan}{'='*60}{reset}")
        print()

        return True

    def speak_sync(self, text: str) -> bool:
        """
        同步播放语音解说（阻塞等待完成）

        Args:
            text (str): 要播放的文本

        Returns:
            bool: 播放成功返回 True
        """
        try:
            if self.tts_engine == "edge" and EDGE_TTS_AVAILABLE:
                return self._speak_edge_sync(text)
            elif PYTTSX3_AVAILABLE:
                return self._speak_pyttsx3_sync(text)
            else:
                logger.warning("⚠️ 没有可用的 TTS 引擎")
                return False

        except Exception as e:
            logger.error(f"❌ TTS 播放异常: {type(e).__name__}: {str(e)}")
            return False

    def _speak_edge_sync(self, text: str) -> bool:
        """使用 edge-tts 同步播放语音"""
        try:
            import tempfile

            logger.info(f"🗣️ [Edge-TTS] 正在生成语音: {text[:30]}...")

            communicate = edge_tts.Communicate(
                text=text,
                voice=DEFAULT_VOICE
            )

            temp_file = tempfile.NamedTemporaryFile(suffix=".mp3", delete=False)
            temp_path = temp_file.name
            temp_file.close()

            # 同步保存音频文件
            import asyncio
            asyncio.run(communicate.save(temp_path))

            logger.debug(f"💾 音频文件已保存: {temp_path}")

            # 播放音频（阻塞等待完成）
            logger.debug("▶️ 正在播放音频...")
            self._play_audio_file(temp_path)

            # 删除临时文件
            try:
                os.remove(temp_path)
                logger.debug("🗑️ 临时音频文件已删除")
            except:
                pass

            logger.info("✅ [Edge-TTS] 语音播放完成")
            return True

        except Exception as e:
            logger.error(f"❌ edge-tts 播放失败: {str(e)}")
            return False

    def _speak_pyttsx3_sync(self, text: str) -> bool:
        """使用 pyttsx3 同步播放语音"""
        try:
            logger.info(f"🗣️ [pyttsx3] 正在播放: {text[:30]}...")

            engine = pyttsx3.init()

            # 尝试设置中文语音
            voices = engine.getProperty('voices')
            for voice in voices:
                if 'zh' in voice.id.lower() or 'chinese' in voice.id.lower():
                    engine.setProperty('voice', voice.id)
                    break

            engine.say(text)
            engine.runAndWait()  # 阻塞等待播放完成
            engine.stop()

            logger.info("✅ [pyttsx3] 语音播放完成")
            return True

        except Exception as e:
            logger.error(f"❌ pyttsx3 播放失败: {str(e)}")
            return False

    def _play_audio_file(self, audio_path: str):
        """播放音频文件（阻塞等待完成）"""
        import platform
        import subprocess

        system = platform.system()

        if system == "Windows":
            cmd = ['powershell', '-c', f'(New-Object Media.SoundPlayer "{audio_path}").PlaySync()']
        elif system == "Darwin":  # macOS
            cmd = ['afplay', audio_path]
        else:  # Linux
            cmd = ['aplay', audio_path]

        subprocess.run(cmd, check=True, capture_output=True)

    def execute_step(self, step: Dict[str, Any], step_number: int, total_steps: int) -> bool:
        """
        执行单个讲解步骤

        Args:
            step (Dict): 步骤数据，包含 text/ra/dec
            step_number (int): 当前步骤编号
            total_steps (int): 总步骤数

        Returns:
            bool: 执行成功返回 True
        """
        text = step.get("text", "")
        ra = step.get("ra", 0.0)
        dec = step.get("dec", 0.0)

        # 打印步骤信息
        cyan = '\033[36m'
        yellow = '\033[33m'
        bold = '\033[1m'
        reset = '\033[0m'

        print()
        print(f"{bold}{cyan}{'─'*60}{reset}")
        print(f"{bold}{yellow}📍 步骤 {step_number}/{total_steps}{reset}")
        print(f"{bold}{cyan}{'─'*60}{reset}")
        logger.info(f"📝 解说词: {text}")
        logger.info(f"🎯 目标坐标: RA={ra:.4f}, DEC={dec:.4f}")
        print(f"{bold}{cyan}{'─'*60}{reset}")
        print()

        # 1. 发送坐标到串口
        send_success = self.send_coordinates(ra, dec)

        # 2. 播放语音解说（同步等待完成）
        speak_success = self.speak_sync(text)

        # 3. 步骤间延迟（让硬件有足够的响应时间）
        if step_number < total_steps:
            delay = self.step_delay
            logger.debug(f"⏱️ 等待 {delay} 秒后进入下一步...")
            time.sleep(delay)

        return send_success and speak_success

    def execute_lesson(self, lesson_script: List[Dict[str, Any]]) -> Dict[str, Any]:
        """
        执行完整的讲课脚本

        Args:
            lesson_script (List[Dict]): 讲课脚本数组

        Returns:
            Dict: 执行结果统计
        """
        total_steps = len(lesson_script)
        completed_steps = 0

        print()
        print(f"\033[1;32m{'='*60}\033[0m")
        print(f"\033[1;32m🎓 开始硬件讲解：共 {total_steps} 个步骤\033[0m")
        print(f"\033[1;32m{'='*60}\033[0m")
        print()

        for idx, step in enumerate(lesson_script, 1):
            try:
                success = self.execute_step(step, idx, total_steps)
                if success:
                    completed_steps += 1
                else:
                    logger.warning(f"⚠️ 步骤 {idx} 执行失败，继续下一步")
            except Exception as e:
                logger.error(f"❌ 步骤 {idx} 执行异常: {type(e).__name__}: {str(e)}")
                continue

        print()
        print(f"\033[1;32m{'='*60}\033[0m")
        print(f"\033[1;32m🎉 硬件讲解完成！{reset} \033[1;36m完成度: {completed_steps}/{total_steps}\033[0m")
        print(f"\033[1;32m{'='*60}\033[0m")
        print()

        return {
            "total_steps": total_steps,
            "completed_steps": completed_steps,
            "success": completed_steps == total_steps
        }

    def cleanup(self):
        """清理资源"""
        if self.serial_conn and self.serial_conn.is_open:
            self.serial_conn.close()
            logger.info("🔌 串口已关闭")


# ============================================
# 主执行函数
# ============================================

def execute_skill(
    lesson_script: List[Dict[str, Any]],
    lecture_md: str,
    ppt_slides: List[Dict[str, Any]]
) -> Dict[str, Any]:
    """
    执行硬件控制技能（课件自动生成版）

    Args:
        lesson_script (List[Dict]): 讲课脚本数组，每个元素包含 text/ra/dec
        lecture_md (str): 长篇讲稿（Markdown 格式）
        ppt_slides (List[Dict]): PPT 页面数组，每页包含 title/bullets

    Returns:
        Dict[str, Any]: 执行结果，格式为 {"status": "success|error", "msg": "消息", ...}
    """
    # 打印执行开始标记
    print()
    print(f"\033[1;36m{'='*60}\033[0m")
    print(f"\033[1;36m🔭 AstroGuide Hardware Skill v2.1 开始执行\033[0m")
    print(f"\033[1;36m{'='*60}\033[0m")
    logger.info(f"📚 接收到讲课脚本 | 共 {len(lesson_script)} 个步骤")
    logger.info(f"📝 接收到讲稿 | {len(lecture_md)} 字")
    logger.info(f"📊 接收到 PPT 大纲 | {len(ppt_slides)} 页")

    try:
        # ========== 阶段 1: 课件生成 ==========
        print()
        print(f"\033[1;35m{'='*60}\033[0m")
        print(f"\033[1;35m📚 [生产力模块] 开始生成课件\033[0m")
        print(f"\033[1;35m{'='*60}\033[0m")
        print()

        # 初始化课件生成器
        generator = CoursewareGenerator()

        # 生成 Markdown 讲稿
        md_file = generator.generate_markdown(lecture_md)

        # 生成 PowerPoint 课件
        pptx_file = generator.generate_powerpoint(ppt_slides)

        # 打印生成结果
        if md_file or pptx_file:
            bold = '\033[1m'
            green = '\033[32m'
            cyan = '\033[36m'
            reset = '\033[0m'

            print()
            print(f"{bold}{green}{'='*60}{reset}")
            print(f"{bold}{green}✅ [生产力模块] 课件生成成功！{reset}")
            print(f"{bold}{green}{'='*60}{reset}")

            if md_file:
                print(f"{cyan}📝 Markdown 讲稿:{reset} {md_file}")
            if pptx_file:
                print(f"{cyan}📊 PowerPoint 课件:{reset} {pptx_file}")

            print(f"{bold}{green}{'='*60}{reset}")
            print()
        else:
            logger.warning("⚠️ 课件生成失败，但将继续执行硬件控制")

        # ========== 阶段 2: 硬件控制 ==========
        # 验证脚本格式
        if not isinstance(lesson_script, list):
            raise ValueError("lesson_script 必须是数组类型")

        if len(lesson_script) == 0:
            raise ValueError("lesson_script 不能为空")

        # 验证每个步骤的格式
        for idx, step in enumerate(lesson_script):
            if not isinstance(step, dict):
                raise ValueError(f"步骤 {idx + 1} 不是对象类型")
            if "text" not in step:
                raise ValueError(f"步骤 {idx + 1} 缺少 text 字段")
            if "ra" not in step:
                raise ValueError(f"步骤 {idx + 1} 缺少 ra 字段")
            if "dec" not in step:
                raise ValueError(f"步骤 {idx + 1} 缺少 dec 字段")

            # 验证坐标范围
            ra = step["ra"]
            dec = step["dec"]
            if not (0 <= ra <= 24):
                raise ValueError(f"步骤 {idx + 1} 的 RA 坐标超出范围 [0, 24]: {ra}")
            if not (-90 <= dec <= 90):
                raise ValueError(f"步骤 {idx + 1} 的 DEC 坐标超出范围 [-90, 90]: {dec}")

        # 初始化硬件控制器
        controller = AstroHardwareController()

        # 执行完整讲课
        result_stats = controller.execute_lesson(lesson_script)

        # 清理资源
        controller.cleanup()

        # 构造返回结果
        files_generated = {}
        if md_file:
            files_generated["markdown_file"] = md_file
        if pptx_file:
            files_generated["pptx_file"] = pptx_file

        if result_stats["success"]:
            result = {
                "status": "success",
                "msg": f"✅ 任务完成！成功生成课件并执行 {result_stats['completed_steps']}/{result_stats['total_steps']} 个硬件步骤",
                "files_generated": files_generated,
                "total_steps": result_stats["total_steps"],
                "completed_steps": result_stats["completed_steps"]
            }
            logger.info(f"🎉 {result['msg']}")
        else:
            result = {
                "status": "success",  # 部分成功也算成功
                "msg": f"⚠️ 任务部分完成：{result_stats['completed_steps']}/{result_stats['total_steps']} 个硬件步骤",
                "files_generated": files_generated,
                "total_steps": result_stats["total_steps"],
                "completed_steps": result_stats["completed_steps"]
            }
            logger.info(f"⚠️ {result['msg']}")

        return result

    except Exception as e:
        error_result = {
            "status": "error",
            "msg": f"❌ 执行失败: {type(e).__name__}: {str(e)}"
        }
        logger.error(error_result["msg"])
        return error_result


# ============================================
# 命令行入口
# ============================================

def main():
    """命令行入口函数"""
    # 解析命令行参数
    parser = argparse.ArgumentParser(
        description='AstroGuide Hardware Skill v2.1 - OpenClaw 硬件控制技能（课件自动生成版）',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
示例用法:
  # 完整调用（包含讲稿和 PPT）
  python run_hardware.py '[{"text":"猎户座介绍","ra":5.5,"dec":7.0}]' '# 猎户座\\n\\n猎户座是...' '[{"title":"猎户座","bullets":["冬季星座","易于辨认"]}]'

环境变量:
  ASTRO_GUIDE_SERIAL_PORT  - 串口设备名称（如 COM3 或 /dev/ttyUSB0）
  ASTRO_GUIDE_TTS_ENGINE   - TTS 引擎（edge 或 pyttsx3，默认 edge）
  ASTRO_GUIDE_STEP_DELAY   - 步骤间延迟（秒，默认 0.5）
  ASTRO_GUIDE_OUTPUT_DIR   - 输出目录（默认 output）
  ASTRO_GUIDE_FORCE_MOCK   - 强制 Mock 模式（设为 "true"）
        """
    )

    parser.add_argument('lesson_script', type=str, help='讲课脚本 JSON 数组')
    parser.add_argument('lecture_md', type=str, help='讲稿内容（Markdown 格式）')
    parser.add_argument('ppt_slides', type=str, help='PPT 页面 JSON 数组')

    args = parser.parse_args()

    # 解析 JSON 数据
    try:
        lesson_script = json.loads(args.lesson_script)
        lecture_md = args.lecture_md  # Markdown 内容保持原样
        ppt_slides = json.loads(args.ppt_slides)
    except json.JSONDecodeError as e:
        print(json.dumps({
            "status": "error",
            "msg": f"❌ JSON 解析失败: {str(e)}"
        }, ensure_ascii=False))
        sys.exit(1)

    # 执行技能
    result = execute_skill(lesson_script, lecture_md, ppt_slides)

    # 输出 JSON 结果（OpenClaw 框架会捕获这个输出）
    print(json.dumps(result, ensure_ascii=False))


if __name__ == "__main__":
    main()
